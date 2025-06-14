\
# core/services/manual_generator_service.py

import os
from typing import List, Dict, Any, Tuple, Optional
import logging
import tempfile
import re

import torch
from PIL import Image
Image.MAX_IMAGE_PIXELS = None # Allow loading large images
from transformers import AutoProcessor

logger = logging.getLogger(__name__)

# Import Qwen VL specific class for multimodal models - based on original implementation
try:
    from transformers import Qwen2_5_VLForConditionalGeneration
    qwen_vl_class = Qwen2_5_VLForConditionalGeneration
    logger.info("Using Qwen2_5_VLForConditionalGeneration for manual generation")
except ImportError:
    try:
        from transformers import Qwen2VLForConditionalGeneration  
        qwen_vl_class = Qwen2VLForConditionalGeneration
        logger.info("Using Qwen2VLForConditionalGeneration as fallback")
    except ImportError:
        logger.error("Could not import Qwen VL classes")
        qwen_vl_class = None

# Try to import PEFT for fine-tuned model support
try:
    from peft import PeftModel, PeftConfig
    PEFT_AVAILABLE = True
    logger.info("PEFT library available for fine-tuned model support")
except ImportError:
    PEFT_AVAILABLE = False
    logger.warning("PEFT library not available - fine-tuned models won't work")

# Add PowerPoint imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from core.config import Settings
# from core.models.manual_generation_document import ManualGenDocument # If passing ORM objects

logger = logging.getLogger(__name__)

# DEVICE = "cuda" if torch.cuda.is_available() else "cpu" # device_map="auto" handles this

class ManualGeneratorService:
    def __init__(self, settings: Settings):
        self.settings = settings
        self.model_name = self.settings.MANUAL_MODEL_NAME
        self.hf_token = self.settings.HUGGING_FACE_TOKEN
        self.image_folder = self.settings.MANUAL_GENERATION_IMAGE_FOLDER
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        
        self.model = None
        self.processor = None
        self._load_model()

    def _load_model(self):
        logger.info(f"Loading manual generator model: {self.model_name}")
        if not self.model_name:
            logger.error("MANUAL_MODEL_NAME is not set in settings. Cannot load model.")
            raise ValueError("MANUAL_MODEL_NAME is not configured.")
        
        if qwen_vl_class is None:
            logger.error("Qwen VL classes are not available. Please install a compatible version of transformers.")
            raise ImportError("Qwen VL classes are not available.")
            
        # Base model for fallback
        base_model_name = "Qwen/Qwen2.5-VL-3B-Instruct"
        model_loaded = False
        
        # Skip fine-tuned model if it's the problematic one and go straight to base model
        if self.model_name == "ARHVNAAG/Manuales_finetuning_generator":
            logger.info("Known problematic model detected, using base model directly")
            self._load_base_model(base_model_name)
            model_loaded = True
        else:
            # Try loading custom model first
            try:
                logger.info(f"Attempting to load model directly from: {self.model_name}")
                self.processor = AutoProcessor.from_pretrained(
                    self.model_name,
                    trust_remote_code=True,
                    token=self.hf_token
                )
                self.model = qwen_vl_class.from_pretrained(
                    self.model_name,
                    torch_dtype=torch.bfloat16,
                    device_map="auto",
                    trust_remote_code=True,
                    token=self.hf_token
                )
                logger.info("✅ Model loaded directly from fine-tuned version")
                model_loaded = True
                
            except Exception as e:
                logger.warning(f"Could not load directly from {self.model_name}: {e}")
                logger.info(f"Attempting to load as PEFT model with base: {base_model_name}")
                
                try:
                    # Import PEFT for adapter loading
                    from peft import PeftModel
                    
                    # Load processor from base model
                    self.processor = AutoProcessor.from_pretrained(
                        base_model_name,
                        trust_remote_code=True,
                        token=self.hf_token
                    )
                    
                    # Load base model first
                    base_model = qwen_vl_class.from_pretrained(
                        base_model_name,
                        torch_dtype=torch.bfloat16,
                        device_map="auto",  
                        trust_remote_code=True,
                        token=self.hf_token
                    )
                    
                    # Apply PEFT adapters
                    self.model = PeftModel.from_pretrained(base_model, self.model_name, token=self.hf_token)
                    logger.info("✅ Base model loaded and PEFT adapters applied")
                    model_loaded = True
                    
                except Exception as e2:
                    logger.error(f"Failed to load as PEFT model: {e2}")
                    logger.warning("Falling back to base model only")
                    self._load_base_model(base_model_name)
                    model_loaded = True
        
        # Configure tokenizer for any loaded model
        self._configure_tokenizer()
        logger.info(f"Manual generator model loaded successfully.")

    def _load_base_model(self, base_model_name: str):
        """Load the base Qwen2.5-VL model as fallback"""
        logger.info(f"Loading base model: {base_model_name}")
        try:
            self.processor = AutoProcessor.from_pretrained(
                base_model_name,
                trust_remote_code=True,
                token=self.hf_token
            )
            self.model = qwen_vl_class.from_pretrained(
                base_model_name,
                torch_dtype=torch.bfloat16,
                device_map="auto",
                trust_remote_code=True,
                token=self.hf_token
            )
            logger.info("✅ Base model loaded successfully")
        except Exception as e:
            logger.error(f"Failed to load base model {base_model_name}: {e}")
            raise RuntimeError(f"Cannot load any model - both custom and base model failed: {e}")

    def _configure_tokenizer(self):
        """Configure tokenizer settings for generation"""
        if not self.processor:
            return
            
        # Configure tokenizer padding (from original implementation)
        if self.processor.tokenizer.pad_token_id is None:
            if self.processor.tokenizer.eos_token_id is not None:
                self.processor.tokenizer.pad_token_id = self.processor.tokenizer.eos_token_id
                logger.info(f"Set pad_token_id to eos_token_id ({self.processor.tokenizer.eos_token_id})")
            else:
                logger.warning("Both pad_token_id and eos_token_id are None - this might cause issues")
        
        # Qwen models prefer left padding for generation
        if hasattr(self.processor.tokenizer, 'padding_side'):
            self.processor.tokenizer.padding_side = "left"
            logger.info("Set tokenizer padding_side to 'left' for generation")

    def _process_vision_info_simple(self, messages):
        """
        Simplified processing for extracting images and text from multimodal messages.
        Based on original implementation.
        """
        images = []
        for msg in messages:
            if "content" in msg:
                for content in msg["content"]:
                    if isinstance(content, dict) and content.get("type") == "image":
                        if "image" in content:
                            # Assume content["image"] is a PIL.Image object
                            images.append(content["image"])
        return images, None

    async def generate_manual_text(
        self, 
        query: str, 
        images_metadata: List[Dict[str, Any]], 
    ) -> Dict[str, str]:
        if not self.model or not self.processor:
            logger.error("Manual generation model not loaded properly.")
            return {"manual_text": "Error: Manual generation model not loaded.", "analysis": ""}
        if not self.image_folder:
            logger.error("MANUAL_GENERATION_IMAGE_FOLDER not configured.")
            return {"manual_text": "Error: Manual generation image folder not configured.", "analysis": ""}
        if not images_metadata:
            return {"manual_text": "No relevant image metadata provided to generate the manual.", "analysis": ""}

        pil_images = []
        image_details_for_prompt = []
        loaded_image_paths = [] # For analysis

        for metadata in images_metadata:
            image_path_relative = metadata.get("image_path")
            if not image_path_relative:
                logger.warning(f"Skipping metadata entry due to missing 'image_path': {metadata.get('id')}")
                continue

            full_image_path = os.path.join(self.image_folder, image_path_relative)
            
            if not os.path.exists(full_image_path):
                logger.warning(f"Image not found at {full_image_path}. Skipping.")
                continue
            
            try:
                image = Image.open(full_image_path).convert("RGB")
                pil_images.append(image)
                loaded_image_paths.append(image_path_relative)
                
                folder_path = os.path.dirname(image_path_relative)
                hierarchy_parts = [part for part in folder_path.split('/') if part]
                hierarchy_str = ' > '.join(hierarchy_parts) if hierarchy_parts else 'Raíz'
                
                image_details_for_prompt.append({
                    "path": image_path_relative,
                    "hierarchy": hierarchy_str,
                    "prompt": metadata.get("prompt", "Sin descripción"),
                    "respuesta": metadata.get("respuesta", "Sin respuesta")
                })
            except Exception as e:
                logger.warning(f"Error loading image {full_image_path}: {e}", exc_info=True)
                continue
        
        if not pil_images:
            return {"manual_text": "No images could be loaded to generate the manual.", "analysis": ""}

        user_text_prompt = f"""Eres un experto en sistemas ERP y creación de manuales de usuario.
Crea un manual detallado y claro que explique: {query}

REGLAS OBLIGATORIAS QUE DEBES SEGUIR:
1. Explica cada paso de forma clara y concisa
2. Describe los elementos de la interfaz que son relevantes
3. Organiza la información en secciones lógicas
4. **CRÍTICO**: Para CADA imagen proporcionada, DEBES incluir al menos una referencia usando exactamente este formato: ![Descripción de la imagen](IMAGE_PATH:ruta/de/la/imagen.png)
5. **OBLIGATORIO**: Usa TODAS las imágenes proporcionadas en tu manual
6. Las referencias IMAGE_PATH: son ESENCIALES para mostrar las imágenes al usuario
7. Aprovecha la estructura jerárquica y el orden de las imágenes para entender el flujo de navegación

Información de las imágenes (ÚSALAS TODAS):
"""
        for i, detail in enumerate(image_details_for_prompt):
            user_text_prompt += f"""
Imagen {i+1}:
- Ruta: {detail['path']}
- Jerarquía: {detail['hierarchy']}
- Prompt: {detail['prompt']}
- Respuesta: {detail['respuesta']}
"""
        
        user_text_prompt += """

FORMATO ESPERADO - SIGUE EXACTAMENTE ESTE PATRÓN:
1. Usa formato markdown para títulos, subtítulos, listas, etc.
2. Para CADA paso del manual, incluye la imagen correspondiente usando: ![Descripción](IMAGE_PATH:ruta/de/la/imagen.png)
3. **OBLIGATORIO**: El manual DEBE incluir referencias IMAGE_PATH: para todas las imágenes proporcionadas
4. Al final del manual, añade una breve nota que explique cómo has utilizado la información de las imágenes

EJEMPLO OBLIGATORIO DE ESTRUCTURA:
## Paso 1: Acceder al módulo
Para acceder al módulo, siga estos pasos:
![Pantalla principal del ERP](IMAGE_PATH:pantalla principal/Catalogos/Impresoras.png)

## Paso 2: Configurar opciones
Configure las opciones como se muestra:
![Configuración de opciones](IMAGE_PATH:otra/ruta/imagen.png)

RECUERDA: DEBES usar el formato IMAGE_PATH: para TODAS las imágenes proporcionadas o el manual será inútil.
"""
        
        generated_text = "Error: Could not generate manual text."
        current_padding_side = self.processor.tokenizer.padding_side

        try:
            logger.info(f"Attempting generation with padding_side: {self.processor.tokenizer.padding_side}")
            messages_for_template = [{"role": "user", "content": []}]
            for img in pil_images:
                 messages_for_template[0]["content"].append({"type": "image", "image": img})
            messages_for_template[0]["content"].append({"type": "text", "text": user_text_prompt})

            # Apply chat template first (following original implementation)
            chat_text_for_model = self.processor.apply_chat_template(
                messages_for_template, tokenize=False, add_generation_prompt=True
            )
            
            try:
                # Method 1: Using process_vision_info_simple (original approach)
                image_inputs, _ = self._process_vision_info_simple(messages_for_template)
                inputs = self.processor(
                    text=[chat_text_for_model], 
                    images=image_inputs, 
                    return_tensors="pt", 
                    padding=True
                ).to(self.model.device)
                
                with torch.no_grad():
                    generated_ids = self.model.generate(
                        **inputs,
                        max_new_tokens=self.settings.MANUAL_GENERATION_MAX_NEW_TOKENS,
                        do_sample=self.settings.MANUAL_GENERATION_DO_SAMPLE,
                        temperature=self.settings.MANUAL_GENERATION_TEMPERATURE,
                        top_p=self.settings.MANUAL_GENERATION_TOP_P
                    )
                
                # Trim input tokens to get only the response
                trimmed_ids = [out[len(inp):] for inp, out in zip(inputs.input_ids, generated_ids)]
                generated_text = self.processor.batch_decode(trimmed_ids, skip_special_tokens=True)[0]
                logger.info("Original method with process_vision_info_simple successful.")
                
            except Exception as e_vision:
                logger.warning(f"Method with process_vision_info_simple failed: {e_vision}")
                logger.info("Trying direct method...")
                
                # Method 2: Direct approach (fallback)
                inputs = self.processor(
                    text=[chat_text_for_model], images=pil_images, return_tensors="pt", padding=True
                ).to(self.model.device)

                with torch.no_grad():
                    output_ids = self.model.generate(
                        **inputs,
                        max_new_tokens=self.settings.MANUAL_GENERATION_MAX_NEW_TOKENS,
                        do_sample=self.settings.MANUAL_GENERATION_DO_SAMPLE,
                        temperature=self.settings.MANUAL_GENERATION_TEMPERATURE,
                        top_p=self.settings.MANUAL_GENERATION_TOP_P
                    )
                
                input_token_len = inputs.input_ids.shape[1]
                generated_ids_only = output_ids[:, input_token_len:]
                generated_text = self.processor.batch_decode(generated_ids_only, skip_special_tokens=True)[0]
                logger.info("Direct method successful.")

        except Exception as e_main:
            logger.error(f"Error with main generation method (padding_side: {current_padding_side}): {e_main}", exc_info=True)
            logger.info("Trying alternative generation method with padding_side='right'...")
            try:
                self.processor.tokenizer.padding_side = "right"
                inputs_alt = self.processor(
                    text=[user_text_prompt], images=pil_images, return_tensors="pt", padding=True
                ).to(self.model.device)

                with torch.no_grad():
                    output_ids_alt = self.model.generate(
                        **inputs_alt,
                        max_new_tokens=self.settings.MANUAL_GENERATION_MAX_NEW_TOKENS,
                        temperature=self.settings.MANUAL_GENERATION_TEMPERATURE,
                        do_sample=self.settings.MANUAL_GENERATION_DO_SAMPLE,
                    )
                input_token_len_alt = inputs_alt.input_ids.shape[1]
                generated_ids_only_alt = output_ids_alt[:, input_token_len_alt:]
                generated_text = self.processor.batch_decode(generated_ids_only_alt, skip_special_tokens=True)[0]
                logger.info("Alternative generation method successful.")
            except Exception as e_alt:
                logger.error(f"Error with alternative generation method: {e_alt}", exc_info=True)
                generated_text = f"Error generating manual with both methods. Main error: {e_main}, Alt error: {e_alt}"
            finally:
                self.processor.tokenizer.padding_side = current_padding_side # Restore
        
        if "Assistant:" in generated_text:
            generated_text = generated_text.split("Assistant:", 1)[1].strip()
        if "<|im_end|>" in generated_text: # Qwen specific
            generated_text = generated_text.split("<|im_end|>")[0].strip()
        if generated_text.startswith(user_text_prompt): # If prompt is repeated
             generated_text = generated_text[len(user_text_prompt):].strip()


        usage_analysis = self._analyze_model_usage(generated_text, loaded_image_paths)
        return {"manual_text": generated_text, "analysis": usage_analysis}
    
    async def generate_powerpoint(
        self, 
        query: str, 
        manual_text: str, 
        image_metadata_list: List[Dict[str, Any]],
        images_base64: Optional[Dict[str, str]] = None
    ) -> str:
        """
        Generate a PowerPoint presentation from manual text and images.
        
        Args:
            query: Original query for the manual
            manual_text: Generated manual text in markdown format
            image_metadata_list: List of image metadata dictionaries
            images_base64: Dictionary of base64 encoded images by path
            
        Returns:
            Path to the generated PowerPoint file
        """
        logger.info(f"Generating PowerPoint presentation for query: '{query}'")
        
        # Create temporary file for the PowerPoint
        temp_dir = tempfile.mkdtemp()
        clean_query = re.sub(r'[^\w\s-]', '', query).strip()
        clean_query = re.sub(r'[-\s]+', '_', clean_query)[:30]
        pptx_filename = f"manual_{clean_query}.pptx"
        output_path = os.path.join(temp_dir, pptx_filename)
        
        try:
            # Try to load template
            template_path = "/root/.ipython/aaaaaaaaaaa/Bnext%20RAGnar/plantilla [Autoguardado].pptx"
            if os.path.exists(template_path):
                prs = Presentation(template_path)
                logger.info("Using PowerPoint template")
            else:
                prs = Presentation()
                logger.info("Creating new PowerPoint without template")
            
            # Slide 1: Title
            if len(prs.slides) == 0 or len(prs.slide_layouts) == 0:
                # Create basic slide if no layouts available
                slide = prs.slides.add_slide(prs.slide_layouts[0] if prs.slide_layouts else None)
            else:
                slide_layout = prs.slide_layouts[0]  # Title layout
                slide = prs.slides.add_slide(slide_layout)
            
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                slide.shapes.title.text = query.replace('?', '').replace('¿', '')
            
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "Manual de Usuario - Sistema ERP"
            
            # Slide 2: Introduction
            if len(prs.slide_layouts) > 1:
                content_layout = prs.slide_layouts[1]  # Content layout
            else:
                content_layout = prs.slide_layouts[0] if prs.slide_layouts else None
                
            intro_slide = prs.slides.add_slide(content_layout)
            if hasattr(intro_slide.shapes, 'title') and intro_slide.shapes.title:
                intro_slide.shapes.title.text = "Introducción"
            
            # Split manual text into sections
            sections = self._parse_manual_sections(manual_text)
            
            # Add content slides
            for section_title, section_content in sections.items():
                content_slide = prs.slides.add_slide(content_layout)
                if hasattr(content_slide.shapes, 'title') and content_slide.shapes.title:
                    content_slide.shapes.title.text = section_title
                
                # Add text content
                if len(content_slide.placeholders) > 1:
                    content_placeholder = content_slide.placeholders[1]
                    content_placeholder.text = section_content[:500]  # Limit text length
            
            # Add images slide
            if image_metadata_list:
                images_slide = prs.slides.add_slide(content_layout)
                if hasattr(images_slide.shapes, 'title') and images_slide.shapes.title:
                    images_slide.shapes.title.text = "Imágenes de Referencia"
                
                temp_files_to_cleanup = []
                
                try:
                    # Add up to 4 images per slide
                    for i, img_meta in enumerate(image_metadata_list[:4]):
                        image_added = False
                        
                        if "image_path" in img_meta:
                            img_path = img_meta["image_path"]
                            
                            # Try to use base64 data first
                            if images_base64 and img_path in images_base64:
                                try:
                                    temp_img_path = self._save_base64_image_to_temp(
                                        images_base64[img_path], img_path
                                    )
                                    if temp_img_path:
                                        temp_files_to_cleanup.append(temp_img_path)
                                        
                                        # Calculate position (2x2 grid)
                                        left = Inches(1 + (i % 2) * 4.5)
                                        top = Inches(2 + (i // 2) * 3)
                                        width = Inches(4)
                                        height = Inches(2.5)
                                        
                                        images_slide.shapes.add_picture(
                                            temp_img_path, left, top, width, height
                                        )
                                        image_added = True
                                        logger.debug(f"Added base64 image {img_path} to PowerPoint")
                                        
                                except Exception as e:
                                    logger.warning(f"Could not add base64 image {img_path} to PowerPoint: {e}")
                            
                            # Fallback to file system image if base64 failed
                            if not image_added:
                                full_img_path = os.path.join(self.image_folder, img_path) if not os.path.isabs(img_path) else img_path
                                
                                if os.path.exists(full_img_path):
                                    try:
                                        # Calculate position (2x2 grid)
                                        left = Inches(1 + (i % 2) * 4.5)
                                        top = Inches(2 + (i // 2) * 3)
                                        width = Inches(4)
                                        height = Inches(2.5)
                                        
                                        images_slide.shapes.add_picture(
                                            full_img_path, left, top, width, height
                                        )
                                        logger.debug(f"Added file system image {img_path} to PowerPoint")
                                    except Exception as e:
                                        logger.warning(f"Could not add file system image {full_img_path} to PowerPoint: {e}")
                
                finally:
                    # Clean up temporary files
                    for temp_file in temp_files_to_cleanup:
                        try:
                            os.unlink(temp_file)
                            logger.debug(f"Cleaned up temporary file: {temp_file}")
                        except Exception as e:
                            logger.warning(f"Could not clean up temporary file {temp_file}: {e}")
            
            # Save PowerPoint
            prs.save(output_path)
            logger.info(f"PowerPoint saved to: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint: {e}", exc_info=True)
            raise
    
    def _parse_manual_sections(self, manual_text: str) -> Dict[str, str]:
        """
        Parse manual text into sections based on markdown headers.
        
        Args:
            manual_text: Manual text in markdown format
            
        Returns:
            Dictionary with section titles as keys and content as values
        """
        sections = {}
        current_section = "Contenido Principal"
        current_content = []
        
        lines = manual_text.split('\n')
        
        for line in lines:
            line = line.strip()
            if line.startswith('#'):
                # Save previous section
                if current_content:
                    sections[current_section] = '\n'.join(current_content)
                    current_content = []
                
                # Start new section
                current_section = line.lstrip('#').strip()
                if not current_section:
                    current_section = "Sección"
            else:
                if line:  # Skip empty lines
                    current_content.append(line)
        
        # Save last section
        if current_content:
            sections[current_section] = '\n'.join(current_content)
        
        # If no sections found, use entire text
        if not sections:
            sections["Contenido Principal"] = manual_text
        
        return sections

    def _analyze_model_usage(self, manual_text: str, image_paths: List[str]) -> str:
        analysis = "\\n\\n## Análisis de uso de información por el modelo (Backend)\\n\\n"
        info_usage_section = False
        if "utiliz" in manual_text.lower() and ("imag" in manual_text.lower() or "imágenes" in manual_text.lower()):
            info_usage_section = True
        
        path_references = set()
        for img_path in image_paths:
            if not img_path: continue
            parts_to_check = [os.path.basename(img_path)]
            parent_folder = os.path.basename(os.path.dirname(img_path))
            if parent_folder: parts_to_check.append(parent_folder)
            
            for part in parts_to_check:
                part_no_ext = os.path.splitext(part)[0]
                if len(part_no_ext) > 3 and part_no_ext.lower() in manual_text.lower():
                    path_references.add(part_no_ext)
                elif len(part) > 3 and part.lower() in manual_text.lower():
                    path_references.add(part)

        if info_usage_section:
            analysis += "✅ El modelo ha incluido una sección explicando cómo utilizó la información proporcionada.\\n"
        else:
            analysis += "⚠️ El modelo no ha incluido una sección explícita sobre cómo utilizó la información. (Esto es opcional)\\n"
        
        if path_references:
            analysis += f"✅ Se encontraron {len(path_references)} referencias a elementos de las rutas/nombres de imágenes.\\n"
            analysis += f"   Referencias encontradas: {', '.join(list(path_references)[:5])}\\n"
        else:
            analysis += "ℹ️ No se encontraron referencias directas a nombres de archivo o carpetas de las imágenes en el texto.\\n"
        
        effectiveness_score = (1 if info_usage_section else 0) + min(1, len(path_references) / 2.0)
        
        if effectiveness_score > 1.0:
            analysis += "\\n**Conclusión**: El modelo parece estar utilizando efectivamente la información proporcionada.\\n"
        elif effectiveness_score > 0.5:
            analysis += "\\n**Conclusión**: El modelo parece estar utilizando parcialmente la información proporcionada.\\n"
        else:
            analysis += "\\n**Conclusión**: El modelo podría no estar referenciando explícitamente las imágenes o sus detalles.\\n"
        return analysis

    def _convert_image_to_base64(self, image_path: str) -> Optional[str]:
        """
        Convert an image file to base64 string.
        
        Args:
            image_path: Full path to the image file
            
        Returns:
            Base64 encoded string or None if error
        """
        try:
            import base64
            
            if not os.path.exists(image_path):
                logger.warning(f"Image not found for base64 conversion: {image_path}")
                return None
                
            with open(image_path, "rb") as image_file:
                encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
                logger.debug(f"Successfully converted image to base64: {image_path}")
                return encoded_string
                
        except Exception as e:
            logger.error(f"Error converting image to base64 {image_path}: {e}")
            return None

    def _save_base64_image_to_temp(self, base64_data: str, image_path: str) -> Optional[str]:
        """
        Save base64 image data to a temporary file for PowerPoint insertion.
        
        Args:
            base64_data: Base64 encoded image string
            image_path: Original image path (used for extension)
            
        Returns:
            Path to temporary image file or None if error
        """
        try:
            import base64
            import tempfile
            
            # Decode base64 data
            image_bytes = base64.b64decode(base64_data)
            
            # Get file extension from original path
            _, ext = os.path.splitext(image_path)
            if not ext:
                ext = '.png'  # Default to PNG
                
            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as temp_file:
                temp_file.write(image_bytes)
                temp_path = temp_file.name
                
            logger.debug(f"Saved base64 image to temporary file: {temp_path}")
            return temp_path
            
        except Exception as e:
            logger.error(f"Error saving base64 image to temp file: {e}")
            return None

