#!/usr/bin/env python3
# filepath: /root/.ipython/generate_manual.py
"""
Script para generar manuales usando el modelo de ColPali para buscar imágenes relevantes 
y el modelo fine-tuned Qwen2.5-VL-3B-Instruct para generar documentación basada en esas imágenes.
"""
import os
import sys
import torch
import pandas as pd
import importlib.util
import io
import json
from sqlalchemy import create_engine, Column, Integer, String, Text, text
from sqlalchemy.orm import declarative_base, sessionmaker
from pgvector.sqlalchemy import Vector
from PIL import Image
from IPython.display import display, Markdown
from tqdm import tqdm
import argparse
from transformers import AutoModelForCausalLM, AutoProcessor
from peft import PeftModel, PeftConfig
from pptx import Presentation
from pptx.util import Inches, Pt
import re

try:
    from transformers import Qwen2_5_VLForConditionalGeneration
except ImportError:
    print("⚠️ No se pudo importar Qwen2_5_VLForConditionalGeneration. El script podría fallar al cargar el modelo VL.")
    Qwen2_5_VLForConditionalGeneration = None

# === CONFIGURACIÓN ===
DB_URL = "postgresql+psycopg2://postgres:postgres@localhost:5432/postgres"
CSV_PATH = "/root/.ipython/erp_metadata.csv"  # Corregido para usar la ruta correcta
IMAGE_FOLDER = "/root/.ipython/ERP_screenshots"
COLPALI_MODEL_NAME = "ARHVNAAG/Bnext"
MANUAL_MODEL_NAME = "ARHVNAAG/Manuales_finetuning_generator"
BASE_MODEL_NAME = "Qwen/Qwen2.5-VL-3B-Instruct"  # Base del modelo fine-tuned
DEVICE = "cuda" if torch.cuda.is_available() else "cpu"
EMBED_DIM = 128
MAX_IMAGES_PER_MANUAL = 5

# Función simplificada para procesar imágenes para modelos Qwen VL
def process_vision_info_simple(messages):
    """
    Procesamiento simplificado para extraer imágenes y texto de mensajes con formato multimodal
    """
    images = []
    for msg in messages:
        if "content" in msg:
            for content in msg["content"]:
                if isinstance(content, dict) and content.get("type") == "image":
                    if "image" in content:
                        # Asumimos que content["image"] es un objeto PIL.Image
                        images.append(content["image"])
    return images, None

# === SETUP DE LA BASE DE DATOS ===
Base = declarative_base()

class Document(Base):
    __tablename__ = "documents"
    id = Column(Integer, primary_key=True)
    image_path = Column(String, unique=True)
    prompt = Column(Text)
    respuesta = Column(Text)
    embedding = Column(Vector(EMBED_DIM))

try:
    engine = create_engine(DB_URL)
    Base.metadata.create_all(engine)
    Session = sessionmaker(bind=engine)
    session = Session()
    DB_AVAILABLE = True
except Exception as e:
    print(f"⚠️ Error conectando a la base de datos: {e}")
    print("⚠️ Continuando sin conexión a la base de datos. Algunas funciones estarán limitadas.")
    DB_AVAILABLE = False

def load_colpali_model():
    """Carga el modelo ColPali para búsqueda de imágenes."""
    from colpali_engine.models import ColPali, ColPaliProcessor

    print("🚀 Cargando modelo ColPali...")
    model = ColPali.from_pretrained(COLPALI_MODEL_NAME, torch_dtype=torch.bfloat16, device_map=DEVICE).eval()
    processor = ColPaliProcessor.from_pretrained(COLPALI_MODEL_NAME)
    return model, processor

def load_manual_generator_model():
    """Carga el modelo fine-tuned para generar manuales."""
    print("🚀 Cargando modelo generador de manuales...")
    
    # Ensure Qwen2_5_VLForConditionalGeneration is available
    if Qwen2_5_VLForConditionalGeneration is None:
        try:
            from transformers import Qwen2_5_VLForConditionalGeneration as LocalQwenVL
        except ImportError:
            raise ImportError("Qwen2_5_VLForConditionalGeneration no se pudo importar. "
                              "Asegúrate de que \'transformers\' esté actualizado y sea compatible.")
        qwen_vl_class = LocalQwenVL
    else:
        qwen_vl_class = Qwen2_5_VLForConditionalGeneration

    try:
        # Intentar cargar directamente desde el modelo finetuneado
        processor = AutoProcessor.from_pretrained(MANUAL_MODEL_NAME, trust_remote_code=True)
        model = qwen_vl_class.from_pretrained(
            MANUAL_MODEL_NAME,
            device_map="auto",
            trust_remote_code=True,
            torch_dtype=torch.bfloat16
        )
        print("✅ Modelo cargado directamente desde la versión finetuneada")
    except Exception as e:
        print(f"⚠️ No se pudo cargar directamente. Intentando cargar como PeftModel: {e}")
        # Cargar el procesador desde el modelo base
        processor = AutoProcessor.from_pretrained(BASE_MODEL_NAME, trust_remote_code=True)
        
        # Cargar el modelo base primero
        model = qwen_vl_class.from_pretrained(
            BASE_MODEL_NAME,
            torch_dtype=torch.bfloat16,
            device_map=DEVICE,
            trust_remote_code=True
        )
        
        # Aplicar los adaptadores LoRA del modelo fine-tuned
        model = PeftModel.from_pretrained(model, MANUAL_MODEL_NAME)
        print("✅ Modelo base cargado y adaptadores PEFT aplicados.")
    
    # Configure tokenizer padding
    if processor.tokenizer.pad_token_id is None:
        processor.tokenizer.pad_token_id = processor.tokenizer.eos_token_id
        print(f"ℹ️ Set processor.tokenizer.pad_token_id to EOS token ID: {processor.tokenizer.eos_token_id}")
    
    # Qwen models often use left-padding for generation
    processor.tokenizer.padding_side = "left"
    print(f"ℹ️ Set processor.tokenizer.padding_side to 'left'")
    
    return model, processor

def find_relevant_images(query, colpali_model, colpali_processor, top_k=3):
    """
    Busca imágenes relevantes en la base de datos usando ColPali.
    
    Args:
        query: String con la consulta del usuario
        colpali_model: Modelo ColPali cargado
        colpali_processor: Procesador del modelo ColPali
        top_k: Número máximo de resultados a devolver
        
    Returns:
        Lista de tuplas (image_path, prompt, respuesta)
    """
    print(f"\n🔎 Buscando imágenes relevantes para: {query}")
    
    
    # Generar vector de consulta
    inputs = colpali_processor.process_queries([query]).to(DEVICE)
    with torch.no_grad():
        output = colpali_model(**inputs)
        query_vector = output.to(torch.float32).cpu().numpy()

        if query_vector.ndim == 2 and query_vector.shape[0] == 1:
            query_vector = query_vector[0]
        elif query_vector.ndim == 3:
            query_vector = query_vector.mean(axis=1).squeeze()

    # Ejecutar búsqueda semántica
    results = session.execute(
        text("""
            SELECT id, image_path, prompt, respuesta
            FROM documents
            ORDER BY embedding <-> CAST(:query_vec AS vector)
            LIMIT :limit
        """),
        {"query_vec": query_vector.tolist(), "limit": top_k}
    ).fetchall()
    
    if not results:
        print("❌ No se encontraron imágenes relevantes.")
        return []
        
    relevant_images = []
    for result in results:
        print(f"🎯 Imagen encontrada: {result.image_path}")
        relevant_images.append((result.image_path, result.prompt, result.respuesta))
        
    return relevant_images

def extract_context_from_path(image_path):
    """
    Extrae información contextual de la ruta de la imagen.
    
    Args:
        image_path: Ruta de la imagen (relativa a IMAGE_FOLDER)
        
    Returns:
        Diccionario con información contextual
    """
    context = {
        "module": None,
        "section": None,
        "subsection": None,
        "function": None,
        "hierarchy_level": 0
    }
    
    path_parts = image_path.split("/")
    filename = os.path.basename(image_path).replace(".png", "").replace(".jpg", "").replace(".jpeg", "")
    
    # Extraer contexto de las partes de la ruta
    if len(path_parts) >= 1:
        if "Catalogos" in path_parts:
            context["module"] = "Catálogos"
            cat_index = path_parts.index("Catalogos")
            if len(path_parts) > cat_index + 1:
                context["section"] = path_parts[cat_index + 1]
                if len(path_parts) > cat_index + 2:
                    context["subsection"] = path_parts[cat_index + 2]
        elif "pantalla principal" in path_parts:
            context["module"] = "Pantalla Principal"
            if "modulo" in filename.lower() or "Modulo" in filename:
                module_name = filename.replace("modulo_", "").replace("Modulo_", "").replace("Modulo ", "")
                context["section"] = f"Módulo {module_name}"
    
    # Determinar nivel jerárquico por profundidad de carpetas
    context["hierarchy_level"] = len(path_parts)
    
    # Intentar extraer función de nombre de archivo
    if "pantalla" in filename.lower():
        context["function"] = "Visualización"
    elif "catalogo" in filename.lower():
        context["function"] = "Administración de catálogos"
    elif "modulo" in filename.lower():
        context["function"] = "Acceso a módulo"
    
    return context

def find_parent_images(image_path, df, max_parent_images=2):
    """
    Encuentra imágenes en carpetas padres para dar contexto adicional.
    Organiza las imágenes por jerarquía para entender la navegación.
    
    Args:
        image_path: Ruta de la imagen principal (relativa a IMAGE_FOLDER)
        df: DataFrame con metadata de las imágenes
        max_parent_images: Máximo número de imágenes padres a incluir
        
    Returns:
        Lista de tuplas (image_path, prompt, respuesta)
    """
    path_parts = image_path.split("/")[:-1] 
    parent_images = []
    
    for i in range(len(path_parts), 0, -1):
        if len(parent_images) >= max_parent_images:
            break
            
        current_folder_to_scan_rel = "/".join(path_parts[:i])
        abs_folder_to_scan = os.path.join(IMAGE_FOLDER, current_folder_to_scan_rel)
        
        if not os.path.isdir(abs_folder_to_scan):
            continue
            
        for file_in_folder in sorted(os.listdir(abs_folder_to_scan)):
            if file_in_folder.lower().endswith((".png", ".jpg", ".jpeg")):
                rel_path_candidate = f"{current_folder_to_scan_rel}/{file_in_folder}"
                
                if os.path.normpath(rel_path_candidate) == os.path.normpath(image_path):
                    continue
                
                is_duplicate = any(os.path.normpath(p_img_path) == os.path.normpath(rel_path_candidate) for p_img_path, _, _ in parent_images)
                if is_duplicate:
                    continue

                # Extraer contexto de la ruta
                context = extract_context_from_path(rel_path_candidate)
                
                metadata_found = False
                if df is not None:
                    metadata_rows = df[df["image_path"] == rel_path_candidate]
                    if not metadata_rows.empty:
                        row = metadata_rows.iloc[0]
                        # Enriquecemos los prompts con información contextual de la ruta
                        prompt = f"{row['funciones_detectadas']} - {context['module']} > {context['section']}" if context['section'] else row['funciones_detectadas']
                        parent_images.append((rel_path_candidate, prompt, row['tipo_pantalla']))
                        metadata_found = True
                        print(f"👨‍👦 Imagen padre encontrada (con metadata): {rel_path_candidate}")
                
                if not metadata_found:
                    # Crear un prompt enriquecido basado en la jerarquía de la ruta
                    folder_description = f"{context['module']} > {context['section']}" if context['section'] else current_folder_to_scan_rel
                    prompt = f"Imagen de contexto en {folder_description}"
                    parent_images.append((rel_path_candidate, prompt, f"Pantalla de {context['function'] if context['function'] else 'contexto general'}"))
                    print(f"👨‍👦 Imagen padre encontrada (sin metadata): {rel_path_candidate}")
                    
                if len(parent_images) >= max_parent_images:
                    break
    
    # Ordenar imágenes por nivel jerárquico para facilitar entendimiento de navegación
    parent_images.sort(key=lambda x: len(x[0].split("/")))
    return parent_images

def analyze_model_usage(manual_text, image_paths):
    """
    Analiza si el modelo realmente está utilizando la información de las imágenes.
    
    Args:
        manual_text: Texto del manual generado
        image_paths: Lista de rutas de imágenes utilizadas
        
    Returns:
        String con análisis de uso de información
    """
    analysis = "\n\n## Análisis de uso de información por el modelo\n\n"
    
    # Verificar si hay una sección que explica cómo se usó la información
    info_usage_section = False
    if "utiliz" in manual_text.lower() and "imag" in manual_text.lower():
        info_usage_section = True
    
    # Buscar si hay referencias a rutas específicas
    path_references = []
    for img_path in image_paths:
        path_parts = img_path.split('/')
        # Buscar menciones de carpetas o nombres de archivos específicos
        for part in path_parts:
            if len(part) > 3 and part.lower() in manual_text.lower():
                path_references.append(part)
    
    # Análisis general
    if info_usage_section:
        analysis += "✅ El modelo ha incluido una sección explicando cómo utilizó la información proporcionada.\n"
    else:
        analysis += "❌ El modelo no ha incluido una sección explícita sobre cómo utilizó la información.\n"
    
    if path_references:
        analysis += f"✅ Se encontraron {len(path_references)} referencias a elementos específicos de las rutas de imágenes.\n"
        analysis += f"   Referencias encontradas: {', '.join(path_references[:5])}\n"
    else:
        analysis += "❌ No se encontraron referencias a elementos específicos de las rutas de imágenes.\n"
    
    # Conclusión
    effectiveness_score = (1 if info_usage_section else 0) + min(1, len(path_references)/2)
    if effectiveness_score > 1:
        analysis += "\n**Conclusión**: El modelo está utilizando efectivamente la información proporcionada.\n"
    elif effectiveness_score > 0:
        analysis += "\n**Conclusión**: El modelo está utilizando parcialmente la información proporcionada.\n"
    else:
        analysis += "\n**Conclusión**: El modelo parece no estar utilizando la información proporcionada con efectividad.\n"
        
    return analysis

def generate_manual(query, images_data, image_dict=None):
    """
    Genera un manual basado en las imágenes encontradas y verifica el uso de la información.
    
    Args:
        query: Consulta original del usuario
        images_data: Lista de tuplas (image_path, prompt, respuesta)
        image_dict: Diccionario opcional con metadata adicional que puede incluir imágenes precargadas
        
    Returns:
        String con el texto del manual generado y análisis de uso de información
    """
    if not images_data:
        return "No se encontraron imágenes relevantes para generar el manual."
        
    # Cargar el modelo generador de manuales
    model, processor = load_manual_generator_model()
    
    # Preparar las imágenes y prompts
    image_objects = []
    prompts = []
    respuestas = []
    
    for img_path, prompt, respuesta in images_data:
        # Primero verificar si la imagen está en el diccionario de metadata
        img = None
        if image_dict and img_path in image_dict and "image" in image_dict[img_path]:
            try:
                img = image_dict[img_path]["image"]
                print(f"✅ Usando imagen precomputada para: {img_path}")
            except Exception as e:
                print(f"⚠️ Error al obtener imagen de metadata: {e}")
        
        # Si no encontramos la imagen en el diccionario, intentar cargarla del archivo
        if img is None:
            # Cargar la imagen
            full_path = os.path.join(IMAGE_FOLDER, img_path)
            if not os.path.exists(full_path):
                print(f"⚠️ Imagen no encontrada: {full_path}")
                continue
                
            try:
                img = Image.open(full_path).convert("RGB")
            except Exception as e:
                print(f"⚠️ Error al cargar la imagen {full_path}: {e}")
                continue
        
        # Si tenemos una imagen válida, añadirla a la lista
        if img is not None:
            image_objects.append(img)
            prompts.append(prompt if prompt else "Sin descripción")
            respuestas.append(respuesta if respuesta else "Sin respuesta")
        else:
            print(f"⚠️ No se pudo obtener una imagen válida para: {img_path}")
    
    if not image_objects:
        return "No se pudieron cargar las imágenes para generar el manual."
    
    # Generar el texto del manual
    print(f"📝 Generando manual con {len(image_objects)} imágenes...")
    
    # Construir el prompt para el modelo
    user_text = f"""Eres un experto en sistemas ERP y creación de manuales de usuario.
Crea un manual detallado y claro que explique: {query}
Utiliza las imágenes proporcionadas como guía y sigue estas instrucciones:
1. Explica cada paso de forma clara y concisa
2. Describe los elementos de la interfaz que son relevantes
3. Organiza la información en secciones lógicas
4. Incluye consejos útiles cuando sea apropiado
5. Aprovecha la estructura jerárquica y el orden de las imágenes para entender el flujo de navegación

Información adicional de las imágenes:
"""

    # Añadir información de las imágenes al prompt con rutas de archivo para mostrar jerarquía/navegación
    for i, ((img_path, prompt, respuesta), img) in enumerate(zip(images_data, image_objects)):
        # Extraer información de la jerarquía de carpetas para entender la navegación del sistema
        folder_path = os.path.dirname(img_path)
        hierarchy = folder_path.split('/')
        
        user_text += f"\nImagen {i+1}:\n- Ruta: {img_path}\n- Jerarquía: {' > '.join(hierarchy) if hierarchy else 'Raíz'}\n- Prompt: {prompt}\n- Respuesta: {respuesta}\n"
        
    # Añadir indicación para verificar si el modelo está utilizando la información proporcionada
    user_text += "\nIMPORTANTE: Al final del manual, añade una breve nota que explique cómo has utilizado la información de las imágenes y su metadata para generar este manual."

    # Crear mensaje con formato de chat
    messages = [{
        "role": "user",
        "content": [
            *[{"type": "image", "image": img} for img in image_objects],
            {"type": "text", "text": user_text}
        ]
    }]
    
    try:
        # Aplicar el formato de chat - esto genera texto no tokenizado
        chat_text = processor.apply_chat_template(
            messages,
            tokenize=False,
            add_generation_prompt=True
        )
        
        # Crear un diccionario de entradas apropiado para el modelo
        try:
            # Procesar imágenes usando nuestra función simplificada
            image_inputs, _ = process_vision_info_simple(messages)
            inputs = processor(
                text=[chat_text], 
                images=image_inputs, 
                return_tensors="pt", 
                padding=True
            ).to(model.device)
            
            with torch.no_grad():
                generated_ids = model.generate(
                    **inputs,
                    max_new_tokens=1024,
                    do_sample=True,
                    temperature=0.7,
                    top_p=0.9
                )
            
            # Recortar los tokens de entrada para quedarnos solo con la respuesta
            trimmed_ids = [out[len(inp):] for inp, out in zip(inputs.input_ids, generated_ids)]
            manual_text = processor.batch_decode(trimmed_ids, skip_special_tokens=True)[0]
        except Exception as e:
            print(f"❌ Error con método principal: {e}")
            
            # Método alternativo más directo
            processor.tokenizer.padding_side = "right"  # Cambiar a paddeo derecho para este método
            inputs = processor(
                text=user_text,
                images=image_objects,
                return_tensors="pt",
                padding=True
            ).to(model.device)
            
            with torch.no_grad():
                output_ids = model.generate(
                    **inputs,
                    max_new_tokens=1024,
                    temperature=0.7,
                    do_sample=True,
                )
            
            # Procesar solo la parte generada (recortar tokens de entrada)
            manual_text = processor.batch_decode(
                output_ids[:, inputs.input_ids.shape[1]:],
                skip_special_tokens=True
            )[0]
        
        # Limpiar la salida
        if "Assistant:" in manual_text:
            manual_text = manual_text.split("Assistant:", 1)[1].strip()
        
        # Analizar si el modelo utilizó la información proporcionada
        image_paths = [path for path, _, _ in images_data]
        usage_analysis = analyze_model_usage(manual_text, image_paths)
        
        # Añadir el análisis al texto del manual
        manual_text_with_analysis = manual_text + usage_analysis
        print(usage_analysis)
        
    except Exception as e:
        print(f"❌ Error al generar el manual: {e}")
        manual_text_with_analysis = f"Error al generar el manual: {e}. Por favor, verifica la compatibilidad del modelo."
    
    return manual_text_with_analysis

def create_pptx_from_manual(query, manual_text, all_images):
    """
    Crea una presentación de PowerPoint basada en el manual generado.
    
    Args:
        query: Título del manual (la consulta original)
        manual_text: Texto generado del manual
        all_images: Lista de tuplas (image_path, prompt, respuesta)
        
    Returns:
        Ruta del archivo PowerPoint creado
    """
    print("\n🔄 Generando presentación PowerPoint...")
    
    # Cargar la plantilla de PowerPoint
    template_path = "/root/.ipython/plantilla [Autoguardado].pptx"
    try:
        prs = Presentation(template_path)
    except Exception as e:
        print(f"⚠️ Error al cargar la plantilla PowerPoint: {e}")
        print("⚠️ Creando una presentación nueva sin plantilla.")
        prs = Presentation()
    
    # Slide 1: Título del manual
    slide_layout = prs.slide_layouts[0]  # Layout de título
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Eliminar '?' para mejorar formato del título
    clean_title = query.replace('?', '')
    title.text = f"{clean_title}"
    subtitle.text = "Manual de Usuario"
    
    # Slide 2: Introducción
    slide_layout = prs.slide_layouts[1]  # Layout de contenido
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Introducción"
    
    content = slide.placeholders[1]
    
    # Extraer la introducción del texto del manual
    # Buscar los primeros párrafos antes de cualquier paso numerado
    intro_paragraphs = []
    for line in manual_text.split('\n'):
        line = line.strip()
        # Si encontramos un paso numerado, terminamos de buscar la introducción
        if re.match(r'^\d+\.', line) or "paso" in line.lower() and re.search(r'\d+', line):
            break
        # Ignorar líneas vacías al principio
        if not line and not intro_paragraphs:
            continue
        # Acumular párrafos no vacíos para la introducción
        if line:
            intro_paragraphs.append(line)
        # Limitar a 3 párrafos como máximo para la diapositiva de introducción
        if len(intro_paragraphs) >= 3:
            break
    
    # Si no encontramos una introducción clara, generar una basada en el query
    if not intro_paragraphs:
        intro_text = f"Este manual explica paso a paso cómo {query.replace('?', '').lower()}."
    else:
        intro_text = "\n".join(intro_paragraphs)
    
    # Verificar que la introducción sea adecuada
    if len(intro_text) < 20:  # Si es muy corta, complementar
        intro_text += f"\n\nEste documento proporciona instrucciones detalladas para {query.replace('?', '').lower()} en el sistema ERP."
    
    content.text = intro_text
    
    # Extraer los pasos del manual
    steps = []
    lines = manual_text.strip().split('\n')
    current_step = ""
    for line in lines:
        # Buscar líneas que indiquen números de pasos
        if re.match(r'^\d+\.', line.strip()):
            if current_step:
                steps.append(current_step)
            current_step = line.strip()
        elif current_step:  # Continuación de un paso existente
            current_step += "\n" + line.strip()
    
    # Añadir el último paso
    if current_step:
        steps.append(current_step)
    
    # Asociar imágenes con pasos (distribuir las imágenes disponibles entre los pasos)
    image_objects = []
    for img_path, _, _ in all_images:
        full_path = os.path.join(IMAGE_FOLDER, img_path)
        if os.path.exists(full_path):
            try:
                image_objects.append((full_path, img_path))
            except Exception as e:
                print(f"⚠️ No se pudo procesar la imagen {img_path}: {e}")
    
    # Crear slides para cada paso
    for i, step in enumerate(steps):
        slide_layout = prs.slide_layouts[2]  # Layout para contenido con imágenes
        slide = prs.slides.add_slide(slide_layout)
        
        # Título del paso
        title = slide.shapes.title
        step_num = step.split('.')[0]
        step_title = f"Paso {step_num}"
        title.text = step_title
        
        # Contenido del paso
        try:
            content = slide.placeholders[1]
            content.text = step[step.find('.')+1:].strip()
        except:
            print(f"⚠️ No se pudo añadir contenido al paso {i+1}")
            tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf.text = step[step.find('.')+1:].strip()
        
        # Añadir imagen si está disponible
        img_index = min(i, len(image_objects)-1) if image_objects else -1
        if img_index >= 0:
            try:
                img_path, img_desc = image_objects[img_index]
                left = Inches(6.5)
                top = Inches(2.5)
                pic_width = Inches(3)
                
                slide.shapes.add_picture(img_path, left, top, width=pic_width)
                
                # Obtener metadata para mostrarla en la diapositiva
                img_data = next((data for data in all_images if data[0] == os.path.relpath(img_path, IMAGE_FOLDER)), None)
                context = extract_context_from_path(os.path.relpath(img_path, IMAGE_FOLDER)) if img_data else {}
                
                # Añadir descripción enriquecida debajo de la imagen
                desc_box = slide.shapes.add_textbox(left, top + Inches(3.2), pic_width, Inches(1.0))
                desc_tf = desc_box.text_frame
                
                # Añadir nombre de archivo
                file_p = desc_tf.add_paragraph()
                file_p.text = os.path.basename(img_desc)
                file_p.font.size = Pt(9)
                file_p.font.italic = True
                
                # Añadir información de contexto si está disponible
                if img_data:
                    path_p = desc_tf.add_paragraph()
                    path_parts = os.path.dirname(img_data[0]).split('/')
                    path_text = ' > '.join([p for p in path_parts if p])
                    path_p.text = f"Ruta: {path_text}" if path_text else "Ruta: Raíz"
                    path_p.font.size = Pt(8)
                    path_p.font.italic = True
                    
                    # Si hay prompt disponible, mostrarlo
                    if img_data[1]:
                        prompt_p = desc_tf.add_paragraph()
                        prompt_text = img_data[1]
                        # Limitar longitud
                        if len(prompt_text) > 50:
                            prompt_text = prompt_text[:47] + "..."
                        prompt_p.text = prompt_text
                        prompt_p.font.size = Pt(8)
            except Exception as e:
                print(f"⚠️ Error al añadir imagen al paso {i+1}: {e}")
    
    # Guardar la presentación
    output_pptx = f"manual_{query.replace('?', '').replace(' ', '_')[:30]}.pptx"
    prs.save(output_pptx)
    print(f"✅ Presentación PowerPoint guardada en: {output_pptx}")
    
    return output_pptx

def main():
    """Función principal que ejecuta todo el proceso."""
    parser = argparse.ArgumentParser(description="Generador de manuales con ColPali y Qwen")
    parser.add_argument("--query", type=str, default="¿Cómo registro una nueva impresora?",
                        help="Consulta para buscar imágenes relevantes")
    parser.add_argument("--format", type=str, choices=["md", "pptx", "both"], default="both",
                        help="Formato de salida: md (Markdown), pptx (PowerPoint), o both (ambos)")
    parser.add_argument("--log-input", type=str, default="",
                        help="Archivo de logs con información de metadata o texto directo con logs")
    args = parser.parse_args()
    
    colpali_model, colpali_processor = load_colpali_model()
    
    relevant_images = find_relevant_images(args.query, colpali_model, colpali_processor)
    
    if not relevant_images:
        print("No se encontraron imágenes para generar el manual.")
        return
    
    # Intentar cargar metadata desde CSV
    df = None
    log_metadata = {}
    try:
        df = pd.read_csv(CSV_PATH)
        print("✅ CSV de metadata cargado correctamente.")
    except Exception as e:
        print(f"⚠️ No se pudo cargar el CSV de metadata: {e}")
        print("⚠️ Intentando recuperar metadata de logs...")
        
        # Si se proporcionó un archivo de logs, cargar desde ahí
        if args.log_input:
            try:
                # Ver si es un archivo o texto directo
                if os.path.exists(args.log_input):
                    with open(args.log_input, 'r', encoding='utf-8') as f:
                        log_content = f.read()
                else:
                    # Asumir que es directamente texto de logs
                    log_content = args.log_input
                
                # Importar dinámicamente el módulo para parsear logs
                import importlib.util
                import sys
                
                # Ruta al módulo parse_metadata.py
                module_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "parse_metadata.py")
                
                # Cargar el módulo dinámicamente
                if os.path.exists(module_path):
                    spec = importlib.util.spec_from_file_location("parse_metadata", module_path)
                    parse_metadata = importlib.util.module_from_spec(spec)
                    sys.modules["parse_metadata"] = parse_metadata
                    spec.loader.exec_module(parse_metadata)
                    
                    # Usar la función para extraer metadata
                    log_metadata = parse_metadata.parse_metadata_from_logs(log_content)
                    print(f"✅ Se extrajo metadata para {len(log_metadata)} imágenes desde los logs.")
                else:
                    print(f"⚠️ No se encontró el módulo parse_metadata.py en {module_path}")
            except Exception as e:
                print(f"⚠️ Error al procesar logs: {e}")
        else:
            print("⚠️ No se proporcionó entrada de logs. Continuando sin metadata adicional.")
    
    # Añadir la capacidad de enriquecer el DataFrame con metadata de logs
    if log_metadata and df is not None:
        # Crear un nuevo DataFrame temporal con la metadata de los logs
        log_df_data = []
        for img_path, metadata in log_metadata.items():
            log_df_data.append({
                "image_path": img_path,
                "funciones_detectadas": metadata.get("funciones_detectadas", ""),
                "tipo_pantalla": metadata.get("tipo_pantalla", "")
            })
        
        if log_df_data:
            log_df = pd.DataFrame(log_df_data)
            # Combinar con el DataFrame original, priorizando el original
            df = pd.concat([df, log_df]).drop_duplicates(subset=["image_path"], keep="first")
            print(f"✅ DataFrame enriquecido con {len(log_df)} registros de metadata de logs.")
    
    # Si no tenemos DataFrame pero tenemos metadata de logs, crear un DataFrame con eso
    elif log_metadata and df is None:
        log_df_data = []
        for img_path, metadata in log_metadata.items():
            log_df_data.append({
                "image_path": img_path,
                "funciones_detectadas": metadata.get("funciones_detectadas", ""),
                "tipo_pantalla": metadata.get("tipo_pantalla", "")
            })
        
        if log_df_data:
            df = pd.DataFrame(log_df_data)
            print(f"✅ Se creó un DataFrame con {len(df)} registros de metadata de logs.")
    
    parent_images = find_parent_images(relevant_images[0][0], df)
    
    all_images = relevant_images + parent_images
    all_images = all_images[:MAX_IMAGES_PER_MANUAL]
    
    # Pasar el diccionario de metadata para usar imágenes precargadas
    manual_text = generate_manual(args.query, all_images, log_metadata)
    
    print("\n==== IMÁGENES UTILIZADAS ====")
    for img_path, prompt, respuesta in all_images:
        full_path = os.path.join(IMAGE_FOLDER, img_path)
        if os.path.exists(full_path):
            print(f"\n🖼️ {full_path}")
            print(f"🧠 Prompt: {prompt}")
            print(f"📝 Respuesta: {respuesta}")
            try:
                img = Image.open(full_path).convert("RGB")
                display(img)
            except Exception as e:
                print(f"⚠️ No se pudo mostrar la imagen: {e}")
    
    print("\n==== MANUAL GENERADO ====")
    display(Markdown(manual_text))
    
    # Guardar en formato Markdown si se solicita
    if args.format in ["md", "both"]:
        output_file = f"manual_{args.query.replace('?', '').replace(' ', '_')[:30]}.md"
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(f"# Manual: {args.query}\n\n")
            f.write(manual_text)
        
        print(f"\n✅ Manual guardado en: {output_file}")
    
    # Generar presentación PowerPoint si se solicita
    if args.format in ["pptx", "both"]:
        pptx_path = create_pptx_from_manual(args.query, manual_text, all_images)

if __name__ == "__main__":
    # Si se detectan errores de metadata en la entrada estándar, usarlos
    import sys
    if not sys.stdin.isatty():  # Verificar si hay entrada desde un pipe
        try:
            # Leer entrada estándar (útil para pipes como 'echo "error logs" | python generate_manual.py')
            stdin_content = sys.stdin.read().strip()
            if stdin_content and ("No se pudo cargar el CSV" in stdin_content or "Imagen padre encontrada" in stdin_content):
                print("📄 Detectada información de metadata en la entrada estándar")
                # Importar módulo para parsear logs
                try:
                    import importlib.util
                    module_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "parse_metadata.py")
                    
                    if os.path.exists(module_path):
                        spec = importlib.util.spec_from_file_location("parse_metadata", module_path)
                        parse_metadata = importlib.util.module_from_spec(spec)
                        sys.modules["parse_metadata"] = parse_metadata
                        spec.loader.exec_module(parse_metadata)
                        
                        # Extraer metadata
                        log_metadata = parse_metadata.parse_metadata_from_logs(stdin_content)
                        print(f"✅ Se extrajo metadata para {len(log_metadata)} imágenes desde stdin.")
                        
                        # Añadir la metadata como argumento
                        sys.argv.extend(["--log-input", stdin_content])
                except Exception as e:
                    print(f"⚠️ Error al procesar metadata de stdin: {e}")
        except Exception as e:
            print(f"⚠️ Error al leer de stdin: {e}")
    
    main()
